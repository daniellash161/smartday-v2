#!/usr/bin/env python3
"""
SmartDay PowerPoint Presentation Generator
Creates a professional product deck with Hebrew RTL support and modern design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Color Palette
COLORS = {
    'navy': RGBColor(15, 36, 64),           # #0F2440
    'teal': RGBColor(63, 175, 163),         # #3FAFA3
    'teal_light': RGBColor(223, 244, 236),  # #DFF4EC
    'blue': RGBColor(125, 183, 232),        # #7DB7E8
    'blue_light': RGBColor(238, 244, 251),  # #EEF4FB
    'beige': RGBColor(244, 232, 210),       # #F4E8D2
    'beige_light': RGBColor(254, 250, 242), # #FEFAF2
    'rose': RGBColor(251, 228, 234),        # #FBE4EA
    'mint': RGBColor(223, 244, 236),        # #DFF4EC
    'lavender': RGBColor(236, 232, 255),    # #ECE8FF
    'gold': RGBColor(216, 162, 58),         # #D8A23A
    'bg_light': RGBColor(245, 247, 250),    # #F5F7FA
    'white': RGBColor(255, 255, 255),
    'text_dark': RGBColor(31, 41, 55),      # #1F2937
    'text_muted': RGBColor(107, 114, 128),  # #6B7280
}

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    """Add colored background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle, meta_text=""):
    """Create cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg_light'])

    # Add centered content
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['navy']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER

    # Meta
    if meta_text:
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
        meta_frame = meta_box.text_frame
        p = meta_frame.paragraphs[0]
        p.text = meta_text
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['text_muted']
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, subtitle, bg_color, accent_color):
    """Create standard content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_color)

    # Add accent bar on left
    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent_color
    left_bar.line.color.rgb = accent_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['navy']
    p.alignment = PP_ALIGN.RIGHT

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_muted']
        p.alignment = PP_ALIGN.RIGHT

    return slide

def add_text_content(slide, content, top, left=0.8, width=8.4, font_size=13, color=None):
    """Add text box to slide"""
    if color is None:
        color = COLORS['text_dark']

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    text_frame = box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = 1.4

def add_bullet_list(slide, items, top, left=1.2, font_size=13):
    """Add bullet list to slide"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(7.8), Inches(4))
    text_frame = box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['text_dark']
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.RIGHT
        p.line_spacing = 1.4

# ============================================================================
# SLIDE 1: COVER
# ============================================================================
add_title_slide("SmartDay", "דשבורד יומי חכם שמרכז מידע, מסנן רעש ומתרגם את היום לפעולות ברורות", "פרויקט מוצר - תכנון, בדיקה ושיפור פרוטוטייפ")

# ============================================================================
# SLIDE 2: PROBLEM
# ============================================================================
slide = add_content_slide("הבעיה: ניהול יום בין יותר מדי מקורות",
                          "מידע יומי מפוזר בין יומן, משימות, מיילים, תשלומים ועדכונים - ללא תיעדוף ברור.",
                          COLORS['rose'], COLORS['teal'])

items = [
    "המשתמש עובר בין מערכות שונות, מזהה לבד מה חשוב, ומתרגם מידע לפעולה באופן ידני"
]
add_bullet_list(slide, items, 2.3)

# Highlight box
highlight = slide.shapes.add_shape(1, Inches(1), Inches(4.5), Inches(8), Inches(0.8))
highlight.fill.solid()
highlight.fill.fore_color.rgb = COLORS['rose']
highlight.line.color.rgb = COLORS['rose']

text_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.55), Inches(7.6), Inches(0.7))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "עומס מידע ללא סדר פעולה"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 3: TARGET AUDIENCE
# ============================================================================
slide = add_content_slide("קהל היעד והצרכים המרכזיים",
                          "התמקדנו במשתמשים עמוסים שמנהלים כמה מקורות מידע במקביל.",
                          COLORS['rose'], COLORS['teal'])

# Left box - Audience
left_box = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
left_box.fill.solid()
left_box.fill.fore_color.rgb = COLORS['white']
left_box.line.color.rgb = COLORS['teal']
left_box.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "👥 קהל יעד"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['navy']
p.alignment = PP_ALIGN.RIGHT

items = ["סטודנטים שעובדים במקביל", "משתמשים עם דדליינים", "אנשים שמנהלים יומן אישי"]
add_bullet_list(slide, items, 3, left=1.1, font_size=12)

# Right box - Needs
right_box = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
right_box.fill.solid()
right_box.fill.fore_color.rgb = COLORS['blue_light']
right_box.line.color.rgb = COLORS['blue']
right_box.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "⚡ צרכים מרכזיים"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['navy']
p.alignment = PP_ALIGN.RIGHT

items = ["להבין מה דחוף היום", "לרכז מידע במקום אחד", "להפוך מיילים למשימות"]
add_bullet_list(slide, items, 3, left=5.4, font_size=12)

# ============================================================================
# SLIDE 4: ASSUMPTIONS
# ============================================================================
slide = add_content_slide("ההשערות הראשוניות בתחילת התהליך",
                          "",
                          COLORS['beige_light'], COLORS['teal'])

items = [
    "💡 ריכוז מידע יפחית עומס - אם המידע יופיע במקום אחד, המשתמש יבין מהר יותר",
    "🎯 משתמשים רוצים פעולה - הערך המרכזי הוא הפיכת מיילים לפעולות",
    "📅 סטודנטים צריכים תכנון קדימה - היום ומחר לא מספיקים",
    "🏦 תובנות פיננסיות צריכות להיות שמרניות - אמינות ושקיפות חשובות"
]
add_bullet_list(slide, items, 2.3, font_size=13)

# ============================================================================
# SLIDE 5: PERSONAS
# ============================================================================
slide = add_content_slide("פרסונות מרכזיות",
                          "שני משתמשים מרכזיים המתמודדים עם אתגר דומה",
                          COLORS['lavender'], COLORS['teal'])

# Persona 1
box1 = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
box1.fill.solid()
box1.fill.fore_color.rgb = COLORS['white']
box1.line.color.rgb = COLORS['teal']
box1.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "👩‍🎓 סטודנטית עובדת, 24"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['navy']
p.alignment = PP_ALIGN.RIGHT

items = ["צרכים: לראות שיעורים ומשימות", "Pain: מידע מפוזר בין יומן ומיילים"]
add_bullet_list(slide, items, 3, left=1.1, font_size=11)

# Persona 2
box2 = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
box2.fill.solid()
box2.fill.fore_color.rgb = COLORS['blue_light']
box2.line.color.rgb = COLORS['blue']
box2.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "👩‍💼 אמא, עובדת, 32"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['navy']
p.alignment = PP_ALIGN.RIGHT

items = ["צרכים: תמונת מצב בבוקר", "Pain: עומס תפקידים"]
add_bullet_list(slide, items, 3, left=5.4, font_size=11)

# ============================================================================
# SLIDE 6: PAIN POINTS
# ============================================================================
slide = add_content_slide("נקודות כאב מרכזיות",
                          "",
                          COLORS['rose'], COLORS['teal'])

items = [
    "📊 עומס מידע - מידע יומי מגיע מיומן, מיילים, משימות, אשראי וחדשות",
    "❓ קושי בתיעדוף - המשתמש לא יודע מה דורש טיפול עכשיו",
    "🔄 מעבר בין אפליקציות - כל מקור מידע בכלי אחר",
    "📧 מיילים בין הכיסאות - מייל חשוב לא הופך למשימה",
    "💾 חוסר רציפות - העדפות לא תמיד נשמרות"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# Highlight
highlight = slide.shapes.add_shape(1, Inches(1), Inches(6), Inches(8), Inches(0.8))
highlight.fill.solid()
highlight.fill.fore_color.rgb = COLORS['rose']
highlight.line.color.rgb = COLORS['rose']

text_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.05), Inches(7.6), Inches(0.7))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "הכאב המרכזי: המידע קיים, אך לא מחובר לפעולה"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 7: JOBS TO BE DONE
# ============================================================================
slide = add_content_slide("Jobs To Be Done",
                          "",
                          COLORS['lavender'], COLORS['teal'])

# JTBD 1
box1 = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
box1.fill.solid()
box1.fill.fore_color.rgb = COLORS['white']
box1.line.color.rgb = COLORS['teal']
box1.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.3))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "📍 JTBD 1: פתוח ויום"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['teal']
p.alignment = PP_ALIGN.RIGHT

items = [
    "משימה: להבין מה דורש טיפול",
    "הצלחה: זיהוי מה דחוף",
    "SmartDay: דשבורד במקום אחד"
]
add_bullet_list(slide, items, 2.8, left=1.1, font_size=11)

# JTBD 2
box2 = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
box2.fill.solid()
box2.fill.fore_color.rgb = COLORS['blue_light']
box2.line.color.rgb = COLORS['blue']
box2.line.width = Pt(2)

title = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.3))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "📍 JTBD 2: הפוך לפעולה"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLORS['blue']
p.alignment = PP_ALIGN.RIGHT

items = [
    "משימה: מידע → פעולה",
    "הצלחה: הוספה מהירה",
    "SmartDay: כפתורי פעולה"
]
add_bullet_list(slide, items, 2.8, left=5.4, font_size=11)

# ============================================================================
# SLIDE 8: USER STORIES
# ============================================================================
slide = add_content_slide("User Stories",
                          "",
                          COLORS['blue_light'], COLORS['blue'])

items = [
    "בתור סטודנטית עובדת, כאשר אני פותחת את היום בבוקר, אני רוצה לראות שיעורים ומשימות במקום אחד",
    "",
    "בתור משתמשת עם הרבה מיילים, כאשר מייל דורש פעולה, אני רוצה להפוך אותו למשימה",
    "",
    "בתור משתמשת עם חיובים קבועים, כאשר צפוי חיוב, אני רוצה לקבל תזכורת מראש",
    "",
    "בתור סטודנטית, כאשר יש אירועים עתידיים, אני רוצה לראות אותם בלוח קדימה"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 9: EXISTING SOLUTIONS
# ============================================================================
slide = add_content_slide("פתרונות קיימים",
                          "כל כלי פותר חלק אחד בלבד",
                          COLORS['blue_light'], COLORS['blue'])

items = [
    "📅 Calendar - טוב: ניהול אירועים | חסר: משימות, מיילים, תשלומים",
    "✓ Task Apps - טוב: משימות | חסר: חיבור אוטומטי למיילים",
    "📧 Gmail - טוב: תקשורת | חסר: הפיכה למשימות",
    "💳 Banking - טוב: פירוט חיובים | חסר: תובנות יומיות בדשבורד"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# Highlight
highlight = slide.shapes.add_shape(1, Inches(1), Inches(6), Inches(8), Inches(0.8))
highlight.fill.solid()
highlight.fill.fore_color.rgb = COLORS['blue']
highlight.line.color.rgb = COLORS['blue']

text_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.05), Inches(7.6), Inches(0.7))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "הפער: אין חיבור בין המידע לפעולה"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 10: MARKET TRENDS
# ============================================================================
slide = add_content_slide("מגמות שוק וטכנולוגיה",
                          "",
                          COLORS['blue_light'], COLORS['blue'])

items = [
    "📱 דשבורדים אישיים - משתמשים רוצים ריכוז במקום אחד",
    "🤖 AI לתיעדוף - כלים עוברים מהצגת מידע להמלצה על פעולה",
    "⚙️ אוטומציה - מיילים, אירועים ועדכונים הופכים לפעולות",
    "🔐 פרטיות ושקיפות - משתמשים רוצים לדעת מאיפה מגיע כל מידע"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# Highlight
highlight = slide.shapes.add_shape(1, Inches(1), Inches(6), Inches(8), Inches(0.8))
highlight.fill.solid()
highlight.fill.fore_color.rgb = COLORS['blue']
highlight.line.color.rgb = COLORS['blue']

text_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.05), Inches(7.6), Inches(0.7))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "המגמות תומכות בSmartDay: ריכוז, סינון, פעולה"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 11: SWOT
# ============================================================================
slide = add_content_slide("ניתוח SWOT",
                          "",
                          COLORS['white'], COLORS['teal'])

# Strengths
box_s = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(2), Inches(4.8))
box_s.fill.solid()
box_s.fill.fore_color.rgb = COLORS['teal_light']
box_s.line.color.rgb = COLORS['teal']

title_s = slide.shapes.add_textbox(Inches(0.95), Inches(2.35), Inches(1.7), Inches(0.3))
tf = title_s.text_frame
p = tf.paragraphs[0]
p.text = "💪 Strengths"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS['teal']
p.alignment = PP_ALIGN.CENTER

items = ["ריכוז מידע", "הפיכה לפעולות", "עבור עובדים"]
add_bullet_list(slide, items, 2.7, left=0.95, font_size=10)

# Weaknesses
box_w = slide.shapes.add_shape(1, Inches(3), Inches(2.2), Inches(2), Inches(4.8))
box_w.fill.solid()
box_w.fill.fore_color.rgb = COLORS['rose']
box_w.line.color.rgb = COLORS['rose']

title_w = slide.shapes.add_textbox(Inches(3.15), Inches(2.35), Inches(1.7), Inches(0.3))
tf = title_w.text_frame
p = tf.paragraphs[0]
p.text = "⚠️ Weaknesses"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS['rose']
p.alignment = PP_ALIGN.CENTER

items = ["תלות API", "ניתוח PDF", "סיווג חד"]
add_bullet_list(slide, items, 2.7, left=3.15, font_size=10)

# Opportunities
box_o = slide.shapes.add_shape(1, Inches(5.2), Inches(2.2), Inches(2), Inches(4.8))
box_o.fill.solid()
box_o.fill.fore_color.rgb = COLORS['blue_light']
box_o.line.color.rgb = COLORS['blue']

title_o = slide.shapes.add_textbox(Inches(5.35), Inches(2.35), Inches(1.7), Inches(0.3))
tf = title_o.text_frame
p = tf.paragraphs[0]
p.text = "🚀 Opportunities"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS['blue']
p.alignment = PP_ALIGN.CENTER

items = ["שילוב AI", "הרחבה", "התאמה"]
add_bullet_list(slide, items, 2.7, left=5.35, font_size=10)

# Threats
box_t = slide.shapes.add_shape(1, Inches(7.4), Inches(2.2), Inches(2), Inches(4.8))
box_t.fill.solid()
box_t.fill.fore_color.rgb = COLORS['beige_light']
box_t.line.color.rgb = COLORS['gold']

title_t = slide.shapes.add_textbox(Inches(7.55), Inches(2.35), Inches(1.7), Inches(0.3))
tf = title_t.text_frame
p = tf.paragraphs[0]
p.text = "🎯 Threats"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLORS['gold']
p.alignment = PP_ALIGN.CENTER

items = ["Google", "פרטיות", "API limits"]
add_bullet_list(slide, items, 2.7, left=7.55, font_size=10)

# ============================================================================
# SLIDE 12: VISION
# ============================================================================
slide = add_content_slide("הפתרון: דשבורד יומי שמחבר מידע לפעולה",
                          "SmartDay מרכז מידע ממקורות שונים ומציג רק את מה שדורש תשומת לב",
                          COLORS['teal_light'], COLORS['teal'])

items = [
    "🎯 יומן + משימות + מיילים + תשלומים + עדכונים = SmartDay",
    "",
    "💡 SmartDay לא מחליף את הכלים הקיימים - הוא מחבר ביניהם"
]
add_bullet_list(slide, items, 2.3, font_size=13)

# Feature tags
features = ["יומן", "משימות", "מיילים", "תשלומים", "עדכונים", "העדפות"]
for i, feature in enumerate(features):
    tag = slide.shapes.add_shape(1, Inches(0.8 + i * 1.5), Inches(4.5), Inches(1.4), Inches(0.5))
    tag.fill.solid()
    tag.fill.fore_color.rgb = COLORS['teal']
    tag.line.color.rgb = COLORS['teal']

    text = slide.shapes.add_textbox(Inches(0.85 + i * 1.5), Inches(4.6), Inches(1.3), Inches(0.3))
    tf = text.text_frame
    p = tf.paragraphs[0]
    p.text = feature
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 13: VALUE PROPOSITION
# ============================================================================
slide = add_content_slide("הצעת הערך הייחודית",
                          "ממידע מפוזר לפעולות ברורות",
                          COLORS['teal_light'], COLORS['teal'])

items = [
    "🎯 תמונת מצב יומית - כל מה שחשוב ליום אחד, במקום אחד",
    "",
    "🌬️ סינון רעש - הצגת מידע רלוונטי בלבד, לא רשימות",
    "",
    "⚡ פעולה מיידית - מייל, חיוב או עדכון → משימה",
    "",
    "👤 חוויה אישית - שמירת העדפות, מצב לילה, תאימות"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 14: COMPETITORS
# ============================================================================
slide = add_content_slide("מתחרים ישירים",
                          "השוואת יכולות",
                          COLORS['white'], COLORS['teal'])

items = [
    "📅 Calendar - Low ריכוז, High אירועים",
    "✓ Task Apps - Medium ריכוז, High משימות",
    "📧 Gmail - Low ריכוז, Low משימות",
    "💳 Banking - Low ריכוז, Medium תובנות",
    "",
    "SmartDay - HIGH ריכוז, HIGH משימות, HIGH תובנות תשלומים"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 15: BUSINESS MODEL CANVAS
# ============================================================================
slide = add_content_slide("Business Model Canvas - שלבים 1-4",
                          "",
                          COLORS['beige_light'], COLORS['gold'])

# BMC boxes
boxes = [
    ("🎯 Customer Segments", ["סטודנטים עובדים", "משתמשים עמוסים", "מנהלי פרויקטים"], Inches(0.8), Inches(2.2)),
    ("💎 Value Proposition", ["דשבורד מרכזי", "תיעדוף מידע", "פעולה ברורה"], Inches(3), Inches(2.2)),
    ("📱 Channels", ["אפליקציית Web", "קהילות סטודנטים", "מוסדות"], Inches(5.2), Inches(2.2)),
    ("🤝 Customer Relationships", ["חוויה אישית", "תזכורות חכמות", "שמירת העדפות"], Inches(7.4), Inches(2.2))
]

for title, items_list, left, top in boxes:
    box = slide.shapes.add_shape(1, left, top, Inches(1.8), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['white']
    box.line.color.rgb = COLORS['gold']
    box.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), Inches(1.6), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['gold']
    p.alignment = PP_ALIGN.CENTER

    add_bullet_list(slide, items_list, float(top.inches) + 0.8, left=float(left.inches) + 0.15, font_size=10)

# ============================================================================
# SLIDE 16: PROTOTYPE GOAL
# ============================================================================
slide = add_content_slide("מטרת הפרוטוטייפ",
                          "",
                          COLORS['mint'], COLORS['teal'])

items = [
    "🧠 בדיקת הבנה - המשתמש מבין מה דחוף היום?",
    "",
    "🔄 בדיקת הפיכה - מידע הופך למשימות?",
    "",
    "✨ בדיקת ערך - חיבור מקורות יוצר ערך?",
    "",
    "📊 בדיקת שימושיות - תובנות משפרות חוויה?"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 17: FEATURES
# ============================================================================
slide = add_content_slide("אפיון פיצ׳רים מרכזיים",
                          "",
                          COLORS['white'], COLORS['teal'])

items = [
    "⭐ Critical: לוח זמנים חכם, מיני לוח שנה, מיילים חשובים",
    "",
    "⭐ Critical: תשלומים עם ניתוח PDF, תזכורות תשלום",
    "",
    "⭐ Critical: שמירת העדפות (מצב לילה, מקור יומן)",
    "",
    "✓ Nice-to-Have: עדכוני בוקר, פינה אישית"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 18: PRIORITIZATION
# ============================================================================
slide = add_content_slide("תעדוף פיצ׳רים לפי סדר עבודה",
                          "",
                          COLORS['white'], COLORS['teal'])

items = [
    "1️⃣ שיפור מודול אשראי - הפיצ׳ר היה לא אמין",
    "",
    "2️⃣ הפרדת סוגי תשלומים - למניעת בלבול",
    "",
    "3️⃣ תזכורות תשלום - הופך מידע לפעולה",
    "",
    "4️⃣ שמירת העדפות - מאפשר שימוש רציף",
    "",
    "5️⃣ מיני לוח שנה - תכנון עתידי"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 19: PROTOTYPE OVERVIEW
# ============================================================================
slide = add_content_slide("הפרוטוטייפ הפונקציונלי",
                          "גרסה עובדת שניתן להתנסות בה",
                          COLORS['teal_light'], COLORS['teal'])

# Screenshot placeholder
screenshot = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(8), Inches(2.5))
screenshot.fill.solid()
screenshot.fill.fore_color.rgb = COLORS['white']
screenshot.line.color.rgb = COLORS['text_muted']
screenshot.line.width = Pt(2)
screenshot.line.dash_style = 2  # dashed

text = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(0.8))
tf = text.text_frame
p = tf.paragraphs[0]
p.text = "[צילום מסך של הדשבורד הראשי]"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text_muted']
p.alignment = PP_ALIGN.CENTER

items = ["⚡ React + TypeScript", "📊 דשבורד יומי", "🔗 חיבורי מידע", "💾 שמירה מקומית"]
add_bullet_list(slide, items, 5.2, font_size=12)

# ============================================================================
# SLIDE 20: DEMO FLOW
# ============================================================================
slide = add_content_slide("הדגמה חיה",
                          "זרימות שמראות מידע שהופך לפעולה",
                          COLORS['white'], COLORS['teal'])

items = [
    "1️⃣ פתיחת הדשבורד וסיכום יומי",
    "",
    "2️⃣ צפייה בלוח זמנים ומיני לוח שנה",
    "",
    "3️⃣ הוספת משימה",
    "",
    "4️⃣ הצגת מייל חשוב והפיכתו למשימה",
    "",
    "5️⃣ פתיחת תובנות תשלומים"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 21: EVOLUTION
# ============================================================================
slide = add_content_slide("התפתחות הפרוטוטייפ",
                          "מ-wireframes לגרסה משופרת",
                          COLORS['white'], COLORS['teal'])

# Before
box_before = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
box_before.fill.solid()
box_before.fill.fore_color.rgb = COLORS['rose']
box_before.line.color.rgb = COLORS['rose']

title_before = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_before.text_frame
p = tf.paragraphs[0]
p.text = "📋 Wireframes"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

items = ["מבנה בסיסי", "חלוקה לאזורים"]
add_bullet_list(slide, items, 3, left=1.1, font_size=12)

# After
box_after = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
box_after.fill.solid()
box_after.fill.fore_color.rgb = COLORS['mint']
box_after.line.color.rgb = COLORS['teal']

title_after = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_after.text_frame
p = tf.paragraphs[0]
p.text = "🚀 משופרת"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['teal']
p.alignment = PP_ALIGN.CENTER

items = ["תובנות חכמות", "לוח שנה קדימה", "עדכונים בחלון צף", "פינה אישית"]
add_bullet_list(slide, items, 3, left=5.4, font_size=12)

# ============================================================================
# SLIDE 22: TESTING ROUNDS
# ============================================================================
slide = add_content_slide("שני סבבי בדיקות",
                          "בדיקה, שיפור, בדיקה חוזרת",
                          COLORS['beige_light'], COLORS['gold'])

# Round 1
box_r1 = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
box_r1.fill.solid()
box_r1.fill.fore_color.rgb = COLORS['white']
box_r1.line.color.rgb = COLORS['gold']

title_r1 = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_r1.text_frame
p = tf.paragraphs[0]
p.text = "🔍 סבב ראשון"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['gold']
p.alignment = PP_ALIGN.CENTER

items = ["בדיקת הדשבורד", "זיהוי בעיות עומס", "בדיקת PDF אשראי"]
add_bullet_list(slide, items, 3, left=1.1, font_size=12)

# Round 2
box_r2 = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
box_r2.fill.solid()
box_r2.fill.fore_color.rgb = COLORS['blue_light']
box_r2.line.color.rgb = COLORS['blue']

title_r2 = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_r2.text_frame
p = tf.paragraphs[0]
p.text = "🎯 סבב שני"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['blue']
p.alignment = PP_ALIGN.CENTER

items = ["בדיקה לאחר שיפורים", "תובנות תשלומים", "מיני לוח שנה", "העדפות"]
add_bullet_list(slide, items, 3, left=5.4, font_size=12)

# ============================================================================
# SLIDE 23: METHODOLOGY
# ============================================================================
slide = add_content_slide("משתתפים, מתודולוגיה, מדדים",
                          "",
                          COLORS['white'], COLORS['teal'])

items = [
    "👥 משתתפים: 6 משתמשים - סטודנטים, עובדים, אמהות",
    "",
    "🔬 מתודולוגיה: בדיקות שמישות, משימות, חשיבה בקול רם",
    "",
    "📊 מדדים:",
    "  • מספר פעולות 'הוסף למשימות'",
    "  • תזכורות שנוספו ליומן",
    "  • מציאת אירוע עתידי",
    "  • אמון בתובנות פיננסיות"
]
add_bullet_list(slide, items, 2.3, font_size=11)

# ============================================================================
# SLIDE 24: METRICS ANALYSIS
# ============================================================================
slide = add_content_slide("ניתוח מדדים",
                          "תוצאות ותובנות",
                          COLORS['white'], COLORS['teal'])

items = [
    "📊 הפיכת מידע: משתמשים הבינו את הערך של פעולה",
    "",
    "🔍 אירוע עתידי: שיפור משמעותי אחרי הוספת לוח שנה",
    "",
    "💳 אמון אשראי: שורת מקור הגביר את האמון",
    "",
    "💾 שמירה: העדפות נשמרות וזמינות לאחר רענון"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# ============================================================================
# SLIDE 25: QUALITATIVE FINDINGS
# ============================================================================
slide = add_content_slide("ממצאים כמותיים ואיכותיים",
                          "שני מדדי הצלחה",
                          COLORS['white'], COLORS['teal'])

# Quantitative
box_q = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(4.1), Inches(4.8))
box_q.fill.solid()
box_q.fill.fore_color.rgb = COLORS['blue_light']
box_q.line.color.rgb = COLORS['blue']

title_q = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_q.text_frame
p = tf.paragraphs[0]
p.text = "📊 Quantitative"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['blue']
p.alignment = PP_ALIGN.CENTER

items = ["פעולות 'הוסף'", "שיפור בלוח שנה", "שמירה אמינה"]
add_bullet_list(slide, items, 3, left=1.1, font_size=12)

# Qualitative
box_ql = slide.shapes.add_shape(1, Inches(5.1), Inches(2.2), Inches(4.1), Inches(4.8))
box_ql.fill.solid()
box_ql.fill.fore_color.rgb = COLORS['lavender']
box_ql.line.color.rgb = COLORS['teal']

title_ql = slide.shapes.add_textbox(Inches(5.3), Inches(2.4), Inches(3.7), Inches(0.4))
tf = title_ql.text_frame
p = tf.paragraphs[0]
p.text = "💬 Qualitative"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['teal']
p.alignment = PP_ALIGN.CENTER

items = ["דשבורד ברור יותר", "אשראי שמרני", "עדכונים מאוזנים"]
add_bullet_list(slide, items, 3, left=5.4, font_size=12)

# ============================================================================
# SLIDE 26: RECOMMENDATIONS
# ============================================================================
slide = add_content_slide("מסקנות והמלצות",
                          "שיפורים לשלב הבא",
                          COLORS['white'], COLORS['teal'])

items = [
    "1️⃣ חזק את עיקרון מידע → פעולה",
    "",
    "2️⃣ דשבורד קומפקטי עם פירוט בחלונות צפים",
    "",
    "3️⃣ הרחב תמיכה בספקים אשראי נוספים",
    "",
    "4️⃣ שפר דיוק תזכורות תשלום",
    "",
    "5️⃣ חבר חדשות דרך Backend יציב",
    "",
    "6️⃣ העמיק התאמה אישית ושמירת העדפות"
]
add_bullet_list(slide, items, 2.3, font_size=11)

# ============================================================================
# SLIDE 27: SUMMARY
# ============================================================================
slide = add_content_slide("סיכום",
                          "SmartDay: מרכז שליטה יומי חכם",
                          COLORS['teal_light'], COLORS['teal'])

items = [
    "✓ משתמשים צריכים פחות עומס ויותר תיעדוף",
    "",
    "✓ דשבורד טוב לא מציג הכול - הוא מציג את מה שחשוב",
    "",
    "✓ פיצ׳רים חכמים חייבים להיות אמינים ושקופים",
    "",
    "✓ התאמה אישית ושמירת העדפות = שימוש חוזר"
]
add_bullet_list(slide, items, 2.3, font_size=12)

# Final highlight
highlight = slide.shapes.add_shape(1, Inches(1), Inches(6), Inches(8), Inches(1))
highlight.fill.solid()
highlight.fill.fore_color.rgb = COLORS['teal']
highlight.line.color.rgb = COLORS['teal']

text_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.15), Inches(7.6), Inches(0.7))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "תמונת מצב אחת • פעולות ברורות • פחות רעש"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# Save presentation
# ============================================================================
output_path = '/Users/daniellashemesh/Desktop/smartday-v2/smartday-presentation.pptx'
prs.save(output_path)
print(f"✅ Presentation saved successfully to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
